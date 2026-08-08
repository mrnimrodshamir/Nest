"""Extracts speaker notes from the built PPTX into a standalone markdown file,
so the notes document can never drift from the deck."""
from pptx import Presentation
from pathlib import Path

HERE = Path(__file__).resolve().parent
prs = Presentation(str(HERE / "NestUp_TelAviv_Municipal_Pilot.pptx"))

TITLES = [
    "Cover", "The emotional truth", "The moment", "Today's fragmented experience",
    "What's actually missing", "Introducing NestUp", "The three layers",
    "Where the city already wins", "The RSVP moment", "Communities that persist",
    "Trust by design", "Why the municipality should care", "Where NestUp fits",
    "The pilot", "Measurement", "Why Tel Aviv", "The ask", "Closing",
]

out = ["# NestUp × Tel Aviv-Yafo — Speaker Notes", "",
       "18 slides. Target run time 15 minutes, leaving 15 for discussion.", "",
       "**Room:** Mayor · Municipality CEO · Head of Innovation · Head of Community", "",
       "**Objective:** secure a six-month pilot in one neighbourhood. Not funding.", "",
       "---", ""]

for i, s in enumerate(prs.slides, start=1):
    title = TITLES[i - 1] if i <= len(TITLES) else f"Slide {i}"
    out.append(f"## {i}. {title}")
    out.append("")
    notes = s.notes_slide.notes_text_frame.text.strip()
    out.append(notes if notes else "_No notes._")
    out.append("")
    out.append("---")
    out.append("")

(HERE / "SPEAKER_NOTES.md").write_text("\n".join(out), encoding="utf-8")
print(f"SPEAKER_NOTES.md  ({len(prs.slides._sldIdLst)} slides)")

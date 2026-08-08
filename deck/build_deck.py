"""
NestUp x Tel Aviv-Yafo -- municipal pilot deck.

Native PPTX: every headline, label and diagram is a real PowerPoint shape or
text box, so the whole deck stays editable. Only photographs and the phone
mockups are images.

Typeface: Plus Jakarta Sans (the product's own face). The TTFs ship alongside
this deck in deck/fonts/ -- install them for exact fidelity, otherwise
PowerPoint substitutes and spacing shifts slightly.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
OUT = HERE / "NestUp_TelAviv_Municipal_Pilot.pptx"

# --- palette (src/theme/colors.ts) -----------------------------------------
BG      = RGBColor(0xFA, 0xF8, 0xF4)
SURFACE = RGBColor(0xFE, 0xFD, 0xFB)
INK     = RGBColor(0x2B, 0x2B, 0x28)
INK_2   = RGBColor(0x72, 0x6F, 0x65)
MUTED   = RGBColor(0x76, 0x73, 0x68)
BORDER  = RGBColor(0xE3, 0xE0, 0xD6)
ROSE    = RGBColor(0xA9, 0x5F, 0x70)
ROSE_T  = RGBColor(0xF8, 0xDE, 0xE3)
ROSE_D  = RGBColor(0x8B, 0x4E, 0x5D)
SAGE    = RGBColor(0x7C, 0x9A, 0x82)
SAGE_T  = RGBColor(0xF1, 0xF5, 0xF1)
SAND    = RGBColor(0xC9, 0xA8, 0x76)
SAND_T  = RGBColor(0xF1, 0xE4, 0xCE)
SKY     = RGBColor(0x8F, 0xB4, 0xC9)
SKY_T   = RGBColor(0xE6, 0xEF, 0xF3)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Plus Jakarta Sans"
W, H = 13.333, 7.5
M = 0.95  # page margin

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)


def slide(bg=BG):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=18, weight=400, color=INK, align=PP_ALIGN.LEFT,
        line=1.25, space_after=0, caps=False):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run(); r.text = ln.upper() if caps else ln
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = weight >= 600
        r.font.color.rgb = color
    return tb


def rect(s, x, y, w, h, fill=None, line=None, radius=True, lw=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def circle(s, x, y, d_, fill):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d_), Inches(d_))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def eyebrow(s, y, text, color=ROSE):
    txt(s, M, y, 8, 0.3, text, size=12, weight=600, color=color, caps=True)


def phone(s, img, x, y, h_in):
    """Device mockup at true 393:852 aspect, with a soft bezel."""
    w_in = h_in * 393 / 852
    bez = 0.055
    frame = rect(s, x - bez, y - bez, w_in + 2 * bez, h_in + 2 * bez, fill=INK, radius=True)
    try:
        frame.adjustments[0] = 0.075
    except Exception:
        pass
    s.shapes.add_picture(str(ASSETS / img), Inches(x), Inches(y), Inches(w_in), Inches(h_in))
    return w_in


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text.strip()


def rule(s, x, y, w, color=BORDER):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(1))
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False


# ===========================================================================
# 1 — COVER
# ===========================================================================
s = slide(BG)
rect(s, 0, 0, W, 0.16, fill=ROSE, radius=False)
txt(s, M, 1.5, 9, 0.4, "A proposal for Tel Aviv-Yafo Municipality", size=15, weight=600, color=ROSE, caps=True)
txt(s, M, 2.15, 10.4, 2.6, "The city already\nbuilt the places.", size=62, weight=700, color=INK, line=1.06)
txt(s, M, 4.35, 9.2, 1.4, "NestUp helps the families in them\nfind each other.", size=27, weight=400, color=INK_2, line=1.35)
rule(s, M, 5.95, 4.2)
txt(s, M, 6.2, 8, 0.4, "Municipal community pilot  ·  2026", size=13, weight=500, color=MUTED)
circle(s, 11.1, 2.3, 1.5, ROSE_T)
circle(s, 10.55, 3.5, 0.9, SAGE_T)
circle(s, 11.75, 4.15, 0.55, SAND_T)
notes(s, """
Open here. Don't rush this slide.

"Tel Aviv has spent decades building extraordinary public space for families.
Playgrounds, libraries, parks, community centres, story hours. The infrastructure
is world class.

What the city cannot build with concrete is the thing that makes a neighbourhood
feel like home — knowing the other parents in it.

That's the gap we'd like to close with you.\"

Set expectation: 15 minutes, then discussion. We're asking for a pilot, not funding.
""")

# ===========================================================================
# 2 — THE EMOTIONAL TRUTH
# ===========================================================================
s = slide(INK)
txt(s, M, 2.5, 11.4, 2.6, "You can live on a street with\nfour hundred families\nand know none of them.",
    size=46, weight=700, color=BG, line=1.22)
rule(s, M, 5.6, 3.0, color=ROSE)
notes(s, """
Slow down. Let it sit for a beat before speaking.

"This is the part of city life nobody plans for. Density does not produce
community. You can be surrounded and still be alone."

Do not add statistics here. The line works because it is recognisable, and every
person in the room has either lived it or watched their children live it.
""")

# ===========================================================================
# 3 — THE MOMENT
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "The moment it matters most")
txt(s, M, 1.7, 6.4, 2.4, "When you become\na parent, your city\ngets smaller.", size=44, weight=700, color=INK, line=1.14)
txt(s, M, 4.6, 5.6, 1.8,
    "Your world narrows to the distance you can push a pram. "
    "The people you used to see disappear into work and evenings. "
    "The park is full of people in exactly your situation — and none of you speak.",
    size=17, weight=400, color=INK_2, line=1.62)
# concentric radius diagram
cx, cy = 10.05, 3.85
for d_, col in [(4.3, SAGE_T), (2.9, ROSE_T), (1.4, ROSE)]:
    circle(s, cx - d_ / 2, cy - d_ / 2, d_, col)
txt(s, cx - 0.7, cy - 0.28, 1.4, 0.6, "500m", size=17, weight=700, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, cx - 2.15, cy + 1.62, 4.3, 0.4, "the pram radius", size=13, weight=500, color=MUTED, align=PP_ALIGN.CENTER)
notes(s, """
"Ask any new parent how far they travel in a normal day. It's about 500 metres.

Inside that circle there are dozens of families on exactly the same schedule,
with exactly the same questions, at exactly the same playground. They are
strangers to each other — not because they want to be, but because there is no
natural way to meet."

This is the slide where the room starts nodding. Pause after it.
""")

# ===========================================================================
# 4 — TODAY'S FRAGMENTED EXPERIENCE
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "Today")
txt(s, M, 1.6, 9, 0.9, "Everything exists. Nothing connects.", size=40, weight=700, color=INK)
cards = [
    ("WhatsApp groups", "Invite-only. You have to already\nknow someone to get in.", ROSE_T, ROSE_D),
    ("Facebook groups", "Noisy, city-wide, and rarely\nabout your street.", SKY_T, RGBColor(0x5C, 0x86, 0xA0)),
    ("Municipal listings", "Excellent programmes that\nparents never see.", SAND_T, RGBColor(0x96, 0x79, 0x4F)),
    ("Word of mouth", "Works beautifully — if you\nalready have a network.", SAGE_T, RGBColor(0x4F, 0x6B, 0x55)),
]
cw, gap = 2.72, 0.28
for i, (t, b, fill, fg) in enumerate(cards):
    x = M + i * (cw + gap)
    rect(s, x, 3.1, cw, 2.5, fill=fill)
    txt(s, x + 0.32, 3.45, cw - 0.6, 0.4, t, size=17, weight=600, color=fg)
    txt(s, x + 0.32, 4.05, cw - 0.55, 1.3, b, size=13, weight=400, color=INK_2, line=1.5)
txt(s, M, 6.05, 11.4, 0.6,
    "Four channels. None of them designed for a parent standing in a park on a Tuesday morning.",
    size=17, weight=500, color=INK)
notes(s, """
Walk left to right, one line each. Don't linger.

"None of these are bad products. They're just not built for this job.

The common failure is the same: they all assume you already have a network.
WhatsApp needs an inviter. Facebook needs you to find the right group.
Municipal listings assume you went looking. Word of mouth assumes you have
someone to hear it from.

The families who need community most — new arrivals, first-time parents, people
who just moved neighbourhood — are exactly the ones these channels fail."

If a Head of Community is present, this is where they will agree out loud.
""")

# ===========================================================================
# 5 — WHAT'S ACTUALLY MISSING
# ===========================================================================
s = slide(ROSE)
txt(s, M, 2.55, 11.4, 2.2, "Parents don't need\nmore information.\nThey need each other.",
    size=48, weight=700, color=WHITE, line=1.2)
rule(s, M, 5.55, 3.0, color=RGBColor(0xF3, 0xB6, 0xC2))
notes(s, """
The pivot of the deck. Say it plainly and stop.

"Every city instinct is to publish more — another listing, another portal,
another newsletter. But the shortage was never information.

A parent doesn't want a directory of story hours. They want to know that Dana
from two streets away will be at the one on Tuesday."
""")

# ===========================================================================
# 6 — INTRODUCING NESTUP
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "NestUp")
txt(s, M, 1.65, 6.0, 2.0, "One map of\nyour neighbourhood.", size=40, weight=700, color=INK, line=1.15)
txt(s, M, 3.75, 5.5, 1.9,
    "Activities other parents are hosting. Places the city has verified. "
    "Events the municipality already runs. In one view, sorted by how close "
    "they are to you right now.",
    size=17, weight=400, color=INK_2, line=1.62)
for i, (dot, lab) in enumerate([(ROSE, "Activities — created by parents"),
                                (SAGE, "Places — verified, family-ready"),
                                (SAND, "Events — from the city")]):
    y = 5.35 + i * 0.46
    circle(s, M, y + 0.05, 0.17, dot)
    txt(s, M + 0.38, y, 5.4, 0.35, lab, size=14, weight=500, color=INK)
phone(s, "screen_discovery.png", 8.55, 0.72, 6.1)
notes(s, """
First look at the product. Hold it up, literally, if you have a device.

"This is the whole idea in one screen. One map. Three kinds of thing on it.

Round markers are activities a parent created. Squares are places the city has
verified as family-ready. Diamonds are municipal events.

A parent doesn't have to know which category they want. They just look at what's
near them."

Note the mixed list below the map — an activity and a place, side by side.
That's deliberate. Real life isn't organised into tabs.
""")

# ===========================================================================
# 7 — THE THREE LAYERS
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "What's on the map")
txt(s, M, 1.6, 9, 0.9, "Three layers of a neighbourhood.", size=40, weight=700, color=INK)
cols = [
    ("Activities", ROSE, ROSE_T,
     "A parent decides to be at the\nplayground at ten and says so.\nOthers join.",
     "21 categories\nCreated by residents"),
    ("Places", SAGE, SAGE_T,
     "Verified, family-ready public\nspace — shade, fencing, toilets,\nstroller access.",
     "44 verified in Tel Aviv\nCurated, not crowdsourced"),
    ("Events", SAND, SAND_T,
     "Story hours, workshops and\nprogrammes the municipality\nalready runs.",
     "From municipal sources\nNow with people attached"),
]
cw, gap = 3.66, 0.36
for i, (t, accent, tint, body, foot) in enumerate(cols):
    x = M + i * (cw + gap)
    rect(s, x, 2.95, cw, 3.35, fill=SURFACE, line=BORDER)
    rect(s, x, 2.95, cw, 0.09, fill=accent, radius=False)
    txt(s, x + 0.38, 3.32, cw - 0.7, 0.4, t, size=22, weight=700, color=INK)
    txt(s, x + 0.38, 3.95, cw - 0.72, 1.3, body, size=14, weight=400, color=INK_2, line=1.55)
    rule(s, x + 0.38, 5.35, cw - 0.76)
    txt(s, x + 0.38, 5.55, cw - 0.72, 0.8, foot, size=12, weight=500, color=MUTED, line=1.5)
txt(s, M, 6.62, 11.4, 0.4,
    "The city supplies two of these three layers already.", size=16, weight=600, color=ROSE)
notes(s, """
"Three layers, and here's the part that matters for this conversation:

The municipality already supplies two of them.

You have the places — 44 of them are already verified and live in the product for
Tel Aviv. You run the events. What's been missing is the third layer: the parents
themselves, and a reason for them to show up at the same time.

NestUp doesn't replace anything the city does. It's the connective tissue between
what you already provide and the families who should be using it."
""")

# ===========================================================================
# 8 — THE CITY DOES THE HARD PART
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "Where the city already wins")
txt(s, M, 1.65, 6.4, 2.5, "A great programme\nwith eight people in it\nis still a great programme.",
    size=34, weight=700, color=INK, line=1.22)
txt(s, M, 4.5, 5.9, 1.9,
    "Municipal story hours, workshops and family programmes are genuinely good. "
    "The hard part was never the programming — it's that a parent has no reason "
    "to believe anyone else will be there.",
    size=17, weight=400, color=INK_2, line=1.62)
rect(s, M, 6.15, 5.9, 0.85, fill=SAND_T)
txt(s, M + 0.35, 6.36, 5.3, 0.5, "Attendance is a community problem, not a marketing problem.",
    size=14, weight=600, color=RGBColor(0x96, 0x79, 0x4F))
phone(s, "screen_event.png", 8.55, 0.72, 6.1)
notes(s, """
"This is a real municipal event — a story hour at Beit Ariela — inside NestUp.

Look at what's been added. Not a better description. Not a nicer photo.
Seven other parents who said they're going.

That single line changes the decision. 'Should I take my toddler across town to a
library event?' is a hard question. 'Should I go to the thing seven other parents
from my area are going to?' is an easy one."

Point at the disclaimer under the button — this is important and we'll come back
to it on the trust slide. NestUp attendance is never presented as registration
with the city.
""")

# ===========================================================================
# 9 — THE RSVP MOMENT
# ===========================================================================
s = slide(SURFACE)
eyebrow(s, 1.15, "The mechanic")
txt(s, M, 1.65, 10.5, 1.0, "One tap turns a listing into a gathering.", size=38, weight=700, color=INK)
steps = [
    ("01", "A city event appears", "Pulled from municipal sources into the same map as everything else."),
    ("02", "A parent says \"I'm going\"", "Visible to other NestUp parents. Never sent to the organiser as a registration."),
    ("03", "Others see people, not text", "\"7 NestUp parents going\" — the strongest reason to attend anything."),
    ("04", "They meet in physical space", "The value lands in the park, the library, the community centre. Not in the app."),
]
cw, gap = 2.72, 0.28
for i, (n, t, b) in enumerate(steps):
    x = M + i * (cw + gap)
    txt(s, x, 3.05, 1.0, 0.5, n, size=13, weight=700, color=ROSE)
    txt(s, x, 3.5, cw - 0.2, 0.8, t, size=18, weight=600, color=INK, line=1.25)
    txt(s, x, 4.5, cw - 0.25, 1.6, b, size=13, weight=400, color=INK_2, line=1.55)
    if i < 3:
        rule(s, x + cw - 0.12, 3.22, 0.22, color=BORDER)
rect(s, M, 6.25, 11.4, 0.72, fill=SAGE_T)
txt(s, M + 0.38, 6.45, 10.8, 0.4,
    "Success is measured in the park, not in screen time. We are not trying to keep anyone in the app.",
    size=15, weight=600, color=RGBColor(0x4F, 0x6B, 0x55))
notes(s, """
Four beats, quickly.

Land hard on step four and the green bar.

"I want to be explicit about something, because it matters to a municipality.

We are not building an engagement product. We do not want your residents spending
more time on their phones. Every mechanic in NestUp is designed to end with people
in physical public space — the space you already paid for.

If our usage numbers went up while park attendance stayed flat, we would consider
the pilot a failure."
""")

# ===========================================================================
# 10 — COMMUNITIES THAT PERSIST
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "Between the meetups")
txt(s, M, 1.65, 6.2, 2.0, "Community isn't\nonly an event.", size=40, weight=700, color=INK, line=1.15)
txt(s, M, 3.7, 5.6, 1.9,
    "Twelve curated community spaces run continuously — sleep, feeding, "
    "development, daycare, local recommendations. Always there between "
    "the meetups, in Hebrew and English.",
    size=17, weight=400, color=INK_2, line=1.62)
tags = ["Baby Sleep", "Breastfeeding", "First-Time Parents", "Daycare & Preschools",
        "Local Recommendations", "Parental Leave"]
x, y = M, 5.35
for t in tags:
    w = 0.135 * len(t) + 0.42
    if x + w > M + 5.9:
        x = M; y += 0.58
    rect(s, x, y, w, 0.44, fill=ROSE_T)
    txt(s, x + 0.21, y + 0.115, w, 0.3, t, size=12, weight=500, color=ROSE_D)
    x += w + 0.16
phone(s, "screen_forums.png", 8.55, 0.72, 6.1)
notes(s, """
"A meetup is a moment. A community is what happens between moments.

Twelve forums, curated by us, not user-generated — which means no empty rooms, no
duplicates, no moderation crisis on day one. A parent awake at 3am with a sleep
question has somewhere to go, and it's full of people from their own city.

All of it works in Hebrew and English, right-to-left included."

If asked why users can't create forums: deliberate for the pilot. Twelve good
rooms beat two hundred empty ones, and it keeps moderation tractable.
""")

# ===========================================================================
# 11 — TRUST BY DESIGN
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "Trust")
txt(s, M, 1.65, 6.4, 2.0, "Built for parents,\nso built carefully.", size=40, weight=700, color=INK, line=1.15)
txt(s, M, 3.7, 5.7, 0.9,
    "A public profile shows enough to feel safe meeting someone, and nothing more.",
    size=17, weight=400, color=INK_2, line=1.55)
shown = ["First name and photo", "General area — never an address", "Mom / Dad / Parent, self-selected",
         "Number of children, not their names"]
hidden = ["Phone number", "Email address", "Home address or coordinates", "Children's birthdates"]
txt(s, M, 4.75, 3.0, 0.3, "Shown", size=12, weight=600, color=SAGE, caps=True)
for i, t in enumerate(shown):
    txt(s, M, 5.15 + i * 0.36, 3.3, 0.3, t, size=13, weight=400, color=INK)
txt(s, M + 3.5, 4.75, 3.0, 0.3, "Never shown", size=12, weight=600, color=MUTED, caps=True)
for i, t in enumerate(hidden):
    txt(s, M + 3.5, 5.15 + i * 0.36, 3.3, 0.3, t, size=13, weight=400, color=MUTED)
phone(s, "screen_profile.png", 8.55, 0.72, 6.1)
notes(s, """
Do not skip this slide with a municipal audience. It is the one that gets you
past legal.

"Parent safety shaped the data model, not the marketing.

Children's exact birthdates never leave the database — ages are coarsened before
they're sent to a device. Exact coordinates are never public. Parent role is
self-selected; we never infer it from a name or a photo.

There are no ratings, no trust scores, no verification badges pretending to be
something they aren't. We'd rather show less and be honest about it."

If asked about identity verification: deliberately deferred. We didn't want to
ship a badge that implies a guarantee we can't make.
""")

# ===========================================================================
# 12 — WHY THE MUNICIPALITY SHOULD CARE
# ===========================================================================
s = slide(INK)
eyebrow(s, 1.15, "The municipal case", color=RGBColor(0xF3, 0xB6, 0xC2))
txt(s, M, 1.6, 10, 0.9, "Four things a city gets.", size=40, weight=700, color=BG)
items = [
    ("Isolation", "Parental isolation is a wellbeing issue the city already cares about — and one it currently has no instrument for."),
    ("Utilisation", "Better attendance at programmes and public spaces the city has already paid to build and run."),
    ("Belonging", "Residents who know their neighbours stay longer, participate more, and ask more of their neighbourhood."),
    ("Insight", "Anonymised, aggregated signal about where families actually gather — and where provision is thin."),
]
cw, gap = 2.72, 0.28
for i, (t, b) in enumerate(items):
    x = M + i * (cw + gap)
    rect(s, x, 2.95, cw, 2.75, fill=RGBColor(0x3D, 0x3B, 0x36))
    rect(s, x, 2.95, cw, 0.07, fill=ROSE, radius=False)
    txt(s, x + 0.32, 3.3, cw - 0.6, 0.4, t, size=19, weight=700, color=BG)
    txt(s, x + 0.32, 3.92, cw - 0.6, 1.7, b, size=13, weight=400,
        color=RGBColor(0xA8, 0xA6, 0x9C), line=1.55)
txt(s, M, 6.15, 11.4, 0.4, "None of these require the city to build or maintain software.",
    size=15, weight=500, color=RGBColor(0xF3, 0xB6, 0xC2))
notes(s, """
This is the slide the CEO and Head of Innovation are waiting for. Slow down.

Isolation — you have mental-health and family-wellbeing objectives already. This
is a lever for them that costs almost nothing.

Utilisation — you've built the parks and you fund the programmes. Attendance is
the return on that investment, and it is currently left to chance.

Belonging — the strongest predictor of whether a family stays in a neighbourhood
is whether they know anyone in it.

Insight — and be careful here — aggregated and anonymised. Where families gather,
where they don't, which areas have thin provision. Never individual movement data.
We will put that in writing.

Then the closing line: the city doesn't build or run anything.
""")

# ===========================================================================
# 13 — WHERE NESTUP FITS
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "Fit")
txt(s, M, 1.6, 10, 0.9, "A layer, not a replacement.", size=40, weight=700, color=INK)
lanes = [
    ("The city provides", SAND, SAND_T,
     "Public space  ·  Verified places  ·  Family programmes  ·  Events calendar"),
    ("NestUp provides", ROSE, ROSE_T,
     "Discovery  ·  Attendance  ·  Community  ·  Local trust"),
    ("Families get", SAGE, SAGE_T,
     "Somewhere to go  ·  Someone to go with  ·  A neighbourhood that feels known"),
]
for i, (t, accent, tint, body) in enumerate(lanes):
    y = 2.95 + i * 1.22
    rect(s, M, y, 11.4, 1.0, fill=tint)
    rect(s, M, y, 0.09, 1.0, fill=accent, radius=False)
    txt(s, M + 0.42, y + 0.2, 3.0, 0.35, t, size=16, weight=700, color=INK)
    txt(s, M + 0.42, y + 0.58, 10.4, 0.35, body, size=14, weight=400, color=INK_2)
    if i < 2:
        txt(s, 6.4, y + 1.0, 0.6, 0.25, "↓", size=15, weight=600, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, M, 6.75, 11.4, 0.4,
    "No procurement of software. No integration burden. No new system for staff to run.",
    size=15, weight=600, color=ROSE)
notes(s, """
"I want to be precise about what we are and aren't asking for.

We are not asking the city to buy, build, host or operate software. There's no new
system for your teams to learn and no data migration.

The city keeps doing exactly what it does today. We sit on top and make it findable
and social."

If the Head of Innovation asks about integration: the events layer already ingests
municipal event data. Deeper integration is possible but is explicitly not required
for the pilot.
""")

# ===========================================================================
# 14 — THE PILOT
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "The ask")
txt(s, M, 1.6, 10, 0.9, "One neighbourhood. Six months.", size=40, weight=700, color=INK)
phases = [
    ("Month 1", "Baseline", "Choose one neighbourhood. Agree success measures and capture a starting baseline with the city's own data."),
    ("Months 2–3", "Seed", "Onboard families through existing municipal touchpoints — tipot chalav, libraries, community centres."),
    ("Months 4–5", "Grow", "Municipal events flow into the product. Parents host their own activities. Forums find their rhythm."),
    ("Month 6", "Review", "Joint review against the agreed measures. Decide together whether to extend, expand or stop."),
]
cw, gap = 2.72, 0.28
for i, (p, t, b) in enumerate(phases):
    x = M + i * (cw + gap)
    rect(s, x, 2.95, cw, 2.9, fill=SURFACE, line=BORDER)
    rect(s, x, 2.95, cw, 0.07, fill=ROSE if i < 3 else SAGE, radius=False)
    txt(s, x + 0.3, 3.28, cw - 0.6, 0.3, p, size=12, weight=600, color=ROSE, caps=True)
    txt(s, x + 0.3, 3.68, cw - 0.6, 0.4, t, size=20, weight=700, color=INK)
    txt(s, x + 0.3, 4.25, cw - 0.6, 1.6, b, size=13, weight=400, color=INK_2, line=1.55)
txt(s, M, 6.25, 11.4, 0.4, "One neighbourhood is deliberate. A pilot that can't fail teaches nothing.",
    size=16, weight=600, color=INK)
notes(s, """
"We're asking for one neighbourhood, not the city. Six months, not a year.

Month one is the most important and the one people skip: we establish a baseline
together, using your data, before we've done anything. Otherwise in month six
we'll both be guessing.

Month six has a real stop option in it. If the measures don't move, we'd rather
you end it than renew out of politeness."

Suggested first neighbourhood: somewhere with strong existing family provision and
a dense young-family population. Florentin, Neve Tzedek, or Old North. Ask them
which they'd choose — it's a good way to get them designing the pilot with you.
""")

# ===========================================================================
# 15 — MEASUREMENT
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "Accountability")
txt(s, M, 1.6, 10.5, 0.9, "How we'll know if it worked.", size=40, weight=700, color=INK)
rows = [
    ("Attendance at municipal family events", "City's own attendance records", "Baseline set in month 1"),
    ("Parent-initiated gatherings per month", "In-product, reported monthly", "Zero at start, by definition"),
    ("Repeat participation", "Share of parents attending more than once", "Baseline set in month 1"),
    ("Self-reported local connection", "Short survey at month 1 and month 6", "Baseline set in month 1"),
    ("Reach into underserved areas", "Distribution of activity across the neighbourhood", "Mapped in month 1"),
]
rect(s, M, 2.9, 11.4, 0.5, fill=INK)
for j, hcol in enumerate(["Measure", "Source", "Starting point"]):
    txt(s, M + 0.35 + j * 3.85, 3.03, 3.6, 0.3, hcol, size=12, weight=600, color=BG, caps=True)
for i, r in enumerate(rows):
    y = 3.4 + i * 0.62
    if i % 2 == 0:
        rect(s, M, y, 11.4, 0.62, fill=SURFACE, radius=False)
    for j, cell in enumerate(r):
        txt(s, M + 0.35 + j * 3.85, y + 0.17, 3.7, 0.35, cell,
            size=13, weight=600 if j == 0 else 400, color=INK if j == 0 else INK_2)
rect(s, M, 6.62, 11.4, 0.52, fill=SAND_T)
txt(s, M + 0.35, 6.76, 10.8, 0.3,
    "We deliberately have no numbers to show you yet. Every baseline is set with the city, in month one.",
    size=13, weight=600, color=RGBColor(0x96, 0x79, 0x4F))
notes(s, """
Be direct about the bottom bar. It will earn you more credibility than a slide of
invented projections would.

"You'll notice there are no impressive numbers on this slide.

That's on purpose. We could have put a forecast up. But any figure we produced
today about Tel Aviv parents would be a guess dressed as evidence, and you'd be
right not to trust it.

Every one of these baselines gets set in month one, with your data, before we
start. Then we're both measuring the same thing."

If pushed for projections: offer to model them jointly once the baseline exists.
Do not invent a number in the room.
""")

# ===========================================================================
# 16 — WHY TEL AVIV
# ===========================================================================
s = slide(ROSE)
eyebrow(s, 1.15, "Why here", color=WHITE)
txt(s, M, 1.6, 10.5, 1.8, "Why Tel Aviv should be first.", size=42, weight=700, color=WHITE)
reasons = [
    ("Density", "Neighbourhoods dense enough that a 500-metre radius contains a real community."),
    ("Provision", "The public space and family programming already exist and are genuinely good."),
    ("Appetite", "A municipality with a track record of testing new things before anyone else."),
    ("Readiness", "The product is live, in Hebrew and English, with 44 Tel Aviv places already verified."),
]
cw, gap = 2.72, 0.28
for i, (t, b) in enumerate(reasons):
    x = M + i * (cw + gap)
    txt(s, x, 3.35, 1.4, 0.4, f"0{i+1}", size=13, weight=700, color=RGBColor(0xF3, 0xB6, 0xC2))
    txt(s, x, 3.75, cw - 0.2, 0.4, t, size=21, weight=700, color=WHITE)
    txt(s, x, 4.35, cw - 0.25, 1.8, b, size=13, weight=400, color=RGBColor(0xF8, 0xDE, 0xE3), line=1.6)
txt(s, M, 6.2, 11.4, 0.5, "The first city to do this gets to shape what it becomes.",
    size=18, weight=600, color=WHITE)
notes(s, """
"Four reasons, and the last one is the practical one: this isn't a concept.

The product is built and running. Hebrew and English, right-to-left. Forty-four
Tel Aviv places already verified and live. It's on TestFlight today.

We're not asking you to fund development. We're asking you to be the first city to
point it at a real neighbourhood."

Then the closing line — first-mover framing works well with an innovation team.
""")

# ===========================================================================
# 17 — THE ASK
# ===========================================================================
s = slide(BG)
eyebrow(s, 1.15, "What we need")
txt(s, M, 1.6, 10, 0.9, "Four things. None of them money.", size=40, weight=700, color=INK)
asks = [
    ("A neighbourhood", "One area, chosen with you, to run the pilot in."),
    ("Introductions", "Access to existing family touchpoints — tipot chalav, libraries, community centres."),
    ("Event data", "The family-facing events you already publish, so they appear where parents are looking."),
    ("A partner", "One person inside the municipality who owns this with us for six months."),
]
for i, (t, b) in enumerate(asks):
    y = 2.95 + i * 0.98
    rect(s, M, y, 11.4, 0.82, fill=SURFACE, line=BORDER)
    circle(s, M + 0.42, y + 0.24, 0.34, ROSE_T)
    txt(s, M + 0.42, y + 0.29, 0.34, 0.3, str(i + 1), size=13, weight=700, color=ROSE_D, align=PP_ALIGN.CENTER)
    txt(s, M + 1.05, y + 0.16, 3.2, 0.35, t, size=17, weight=600, color=INK)
    txt(s, M + 4.3, y + 0.19, 6.9, 0.35, b, size=14, weight=400, color=INK_2)
txt(s, M, 6.95, 11.4, 0.4, "We bring the product, the team and the cost.", size=16, weight=600, color=ROSE)
notes(s, """
"Four asks, and I want to be clear that none of them is a budget line.

A neighbourhood. Introductions to the places parents already go. The event data
you already publish. And one person who owns this with us — that's the one that
actually determines whether a pilot works.

We bring the product, the team and the cost."

Then stop talking. Let them respond.
""")

# ===========================================================================
# 18 — CLOSING
# ===========================================================================
s = slide(INK)
txt(s, M, 2.4, 11.4, 2.6, "Cities are built\nfrom buildings.\n\nCommunities are built\nfrom Tuesday mornings.",
    size=42, weight=700, color=BG, line=1.24)
rule(s, M, 6.3, 3.0, color=ROSE)
txt(s, M, 6.55, 8, 0.4, "NestUp  ·  nestup.app", size=13, weight=500, color=RGBColor(0xA8, 0xA6, 0x9C))
notes(s, """
Final slide. Don't add anything to it.

"Tel Aviv already built the buildings. The parks, the libraries, the community
centres — they're there, they're good, and they're often half empty on a Tuesday
morning while a hundred parents sit alone within walking distance of them.

We'd like to help you fill them. One neighbourhood, six months, and an honest
measurement at the end."

Then: "What would you need to see to say yes?"

That question, not a summary, is how you end.
""")

prs.save(OUT)
print(f"saved {OUT.name}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

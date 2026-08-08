"""
Renders PNG + PDF straight from the built PPTX.

Reads back the actual shape geometry, fills and text runs from the .pptx rather
than re-declaring the layout, so an export can never drift from the file the
client opens. Covers exactly the primitives this deck uses: rectangles,
rounded rectangles, ovals, text boxes and pictures.
"""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
from pathlib import Path
import io

HERE = Path(__file__).resolve().parent
FONTS = HERE.parent / "node_modules/@expo-google-fonts/plus-jakarta-sans"
PPTX = HERE / "NestUp_TelAviv_Municipal_Pilot.pptx"
PNG_DIR = HERE / "slides_png"
PNG_DIR.mkdir(exist_ok=True)
PDF = HERE / "NestUp_TelAviv_Municipal_Pilot.pdf"

DPI = 192                       # 13.333in -> 2560px
EMU_IN = 914400
PX = DPI / EMU_IN               # emu -> px


def px(v):
    return int(round(v * PX))


_cache = {}


def font(size_pt, bold):
    key = (round(size_pt, 1), bold)
    if key not in _cache:
        name = "PlusJakartaSans_700Bold" if bold else "PlusJakartaSans_400Regular"
        _cache[key] = ImageFont.truetype(str(FONTS / f"{name}.ttf"), max(1, int(size_pt * DPI / 72)))
    return _cache[key]


def rgb(color_fmt):
    try:
        c = color_fmt.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return None


def shape_fill(shp):
    try:
        if shp.fill.type is not None and shp.fill.type == 1:  # solid
            return rgb(shp.fill.fore_color)
    except Exception:
        pass
    return None


def shape_line(shp):
    try:
        c = shp.line.color
        if c and c.type is not None:
            return rgb(c)
    except Exception:
        pass
    return None


def wrap(draw, text, fnt, max_w):
    """Greedy word wrap to the shape's width, matching PowerPoint closely enough."""
    out = []
    for para in text.split("\n"):
        if not para:
            out.append("")
            continue
        line, words = "", para.split(" ")
        for w in words:
            trial = f"{line} {w}".strip()
            if draw.textlength(trial, font=fnt) <= max_w or not line:
                line = trial
            else:
                out.append(line)
                line = w
        out.append(line)
    return out


def render(slide, idx):
    W = px(prs.slide_width)
    H = px(prs.slide_height)
    im = Image.new("RGB", (W, H), (250, 248, 244))
    d = ImageDraw.Draw(im)

    for shp in slide.shapes:
        try:
            L, T = px(shp.left), px(shp.top)
            Wd, Ht = px(shp.width), px(shp.height)
        except Exception:
            continue

        # --- picture ---
        if shp.shape_type == 13 or getattr(shp, "image", None) is not None:
            try:
                pic = Image.open(io.BytesIO(shp.image.blob)).convert("RGBA")
                pic = pic.resize((max(1, Wd), max(1, Ht)), Image.LANCZOS)
                im.paste(pic, (L, T), pic)
                continue
            except Exception:
                pass

        # --- autoshape fill/outline ---
        fill = shape_fill(shp)
        line = shape_line(shp)
        st = str(getattr(shp, "shape_type", ""))
        name = (shp.name or "").lower()
        is_oval = "OVAL" in st or "oval" in name
        is_round = "ROUNDED" in st or "rounded" in name

        if fill or line:
            box = [L, T, L + Wd, T + Ht]
            if is_oval:
                d.ellipse(box, fill=fill, outline=line, width=max(1, px(12700)))
            elif is_round:
                r = int(min(Wd, Ht) * 0.16)
                d.rounded_rectangle(box, radius=r, fill=fill, outline=line, width=max(1, px(12700)))
            else:
                d.rectangle(box, fill=fill, outline=line, width=max(1, px(12700)))

        # --- text ---
        if not shp.has_text_frame:
            continue
        tf = shp.text_frame
        y = T
        for p in tf.paragraphs:
            runs = [r for r in p.runs if r.text]
            if not runs:
                y += int(font(14, False).size * 1.2)
                continue
            r0 = runs[0]
            size = r0.font.size.pt if r0.font.size else 18
            bold = bool(r0.font.bold)
            col = rgb(r0.font.color) or (43, 43, 40)
            fnt = font(size, bold)
            text = "".join(r.text for r in runs)
            ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.25
            lh = int(size * DPI / 72 * ls)
            for ln in wrap(d, text, fnt, Wd):
                if p.alignment == PP_ALIGN.CENTER:
                    x = L + (Wd - d.textlength(ln, font=fnt)) / 2
                elif p.alignment == PP_ALIGN.RIGHT:
                    x = L + Wd - d.textlength(ln, font=fnt)
                else:
                    x = L
                d.text((x, y), ln, font=fnt, fill=col)
                y += lh
            if p.space_after:
                y += int(p.space_after.pt * DPI / 72)
    return im


if __name__ == "__main__":
    prs = Presentation(str(PPTX))
    pages = []
    for i, s in enumerate(prs.slides, start=1):
        img = render(s, i)
        img.save(PNG_DIR / f"slide_{i:02d}.png")
        pages.append(img.convert("RGB"))
    pages[0].save(PDF, save_all=True, append_images=pages[1:], resolution=DPI)
    print(f"{len(pages)} slides -> slides_png/ and {PDF.name}")
